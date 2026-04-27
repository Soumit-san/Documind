"""
Document Parser Service
Extracts text from PDF, DOCX, PPTX, and plain text files.
Uses PyMuPDF, python-docx, and python-pptx.
"""

import io
import logging

logger = logging.getLogger("documind.parser")


def parse_document(contents: bytes, filename: str) -> dict:
    """
    Parse raw file bytes into plain text + metadata.

    Returns:
        {
            "text": str,           # full extracted text
            "pages": list[str],    # text per page (if applicable)
            "metadata": {
                "filename": str,
                "page_count": int,
                "file_type": str,
            }
        }
    """
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    if ext == ".pdf":
        return _parse_pdf(contents, filename)
    elif ext in (".docx", ".doc", ".odt"):
        return _parse_docx(contents, filename)
    elif ext in (".pptx", ".ppt", ".odp"):
        return _parse_pptx(contents, filename)
    elif ext in (".txt", ".md", ".csv", ".json", ".yaml"):
        return _parse_text(contents, filename, ext)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def _parse_pdf(contents: bytes, filename: str) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=contents, filetype="pdf")
    pages = []
    for page in doc:
        text = page.get_text("text")
        pages.append(text)
    doc.close()

    full_text = "\n\n".join(pages)
    logger.info("Parsed PDF '%s': %d pages, %d chars", filename, len(pages), len(full_text))

    return {
        "text": full_text,
        "pages": pages,
        "metadata": {
            "filename": filename,
            "page_count": len(pages),
            "file_type": "pdf",
        },
    }


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def _parse_docx(contents: bytes, filename: str) -> dict:
    from docx import Document

    doc = Document(io.BytesIO(contents))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)

    logger.info("Parsed DOCX '%s': %d paragraphs, %d chars", filename, len(paragraphs), len(full_text))

    return {
        "text": full_text,
        "pages": [full_text],  # DOCX doesn't have page concept natively
        "metadata": {
            "filename": filename,
            "page_count": 1,
            "file_type": "docx",
        },
    }


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def _parse_pptx(contents: bytes, filename: str) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(contents))
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        slides_text.append("\n".join(texts))

    full_text = "\n\n".join(slides_text)
    logger.info("Parsed PPTX '%s': %d slides, %d chars", filename, len(slides_text), len(full_text))

    return {
        "text": full_text,
        "pages": slides_text,
        "metadata": {
            "filename": filename,
            "page_count": len(slides_text),
            "file_type": "pptx",
        },
    }


# ---------------------------------------------------------------------------
# Plain Text
# ---------------------------------------------------------------------------
def _parse_text(contents: bytes, filename: str, ext: str) -> dict:
    text = contents.decode("utf-8", errors="replace")

    logger.info("Parsed text file '%s': %d chars", filename, len(text))

    return {
        "text": text,
        "pages": [text],
        "metadata": {
            "filename": filename,
            "page_count": 1,
            "file_type": ext.lstrip("."),
        },
    }
